from pathlib import Path
from collections import defaultdict

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml


MAX_SCAN_ROWS = 400
MAX_TABLE_HEIGHT_CM = 14.8
ROW_HEIGHT_HEADER_CM = 1.92
ROW_HEIGHT_DATA_CM = 0.7
MAIN_HEIGHT_LAST_CM = 12.1

COL1_WIDTH_CM = 15.0
COL2_WIDTH_CM = 2.9
COL3_WIDTH_CM = 2.9
COL4_WIDTH_CM = 3.5


def load_book(path: Path):
    return load_workbook(path, data_only=True)


def get_skip_columns_flag(wb):
    try:
        ws = wb["Расчет стоимости"]
    except KeyError:
        return False
    return bool(ws["I1"].value)


def get_headers_sheet(wb):
    try:
        return wb.worksheets[10]
    except IndexError:
        raise RuntimeError("Нужен минимум 11 листов (служебный лист с заголовками).")


def get_header_text(wb, ws, sheet_index: int):
    ws1 = wb.worksheets[0]
    ws11 = get_headers_sheet(wb)

    c1 = ws["C1"].value or ""
    g_val = ws1[f"G{sheet_index - 2}"].value or ""
    a2 = ws11["A2"].value or ""

    parts = [str(c1)]
    if g_val:
        parts.append(str(g_val))
    if a2:
        parts.append(str(a2))
    if len(parts) == 1:
        return parts[0]
    return ", ".join(parts[:-1]) + " " + parts[-1]


def collect_rows_for_sheet(wb, ws, skip_columns: bool):
    ws11 = get_headers_sheet(wb)
    hdr_w = str(ws11["A1"].value or "")
    hdr_p = str(ws11["B1"].value or "")
    hdr_g = str(ws11["C1"].value or "")

    rows_raw = []
    last_row = min(ws.max_row, MAX_SCAN_ROWS)

    for i in range(2, last_row + 1):
        cell_E = ws[f"E{i}"].value
        if cell_E is None or cell_E == 0:
            continue

        cat_cell = str(ws[f"B{i}"].value or "")
        name_cell = str(ws[f"C{i}"].value or "")
        d_val = ws[f"D{i}"].value
        e_val = ws[f"E{i}"].value
        f_val = ws[f"F{i}"].value

        if ("Категория блюд" in cat_cell) or ("Наименован" in name_cell) \
                or (str(d_val) == hdr_w) or (str(e_val) == hdr_p) or (str(f_val) == hdr_g):
            continue

        category = cat_cell
        name = name_cell
        weight = d_val if not skip_columns else None
        portions = e_val if not skip_columns else None
        g_per_person = f_val

        rows_raw.append([category, name, weight, portions, g_per_person])

    return rows_raw


def get_category_order_from_ae(wb):
    ws_cat = wb.worksheets[2]
    order = []
    row = 3
    while True:
        val = ws_cat[f"AE{row}"].value
        if val in (None, ""):
            break
        val_str = str(val)
        if val_str not in order:
            order.append(val_str)
        row += 1
    return order


def build_master_rows_and_totals(wb, rows_raw, skip_columns: bool):
    ws11 = get_headers_sheet(wb)

    valid_categories = set()
    for row in ws11["A8:A12"]:
        for cell in row:
            if cell.value not in (None, ""):
                valid_categories.add(str(cell.value))

    total_food = 0.0
    total_liquid = 0.0

    for row in rows_raw:
        cat_name = str(row[0] or "")
        val = row[4]
        if isinstance(val, (int, float)):
            if cat_name in valid_categories:
                total_liquid += float(val)
            else:
                total_food += float(val)

    cat_to_rows: dict[str, list[list]] = defaultdict(list)
    for category, name, weight, portions, gpp in rows_raw:
        cat = str(category or "")
        if not cat:
            continue
        cat_to_rows[cat].append([category, name, weight, portions, gpp])

    categories_in_data = list(cat_to_rows.keys())
    category_order_from_ae = get_category_order_from_ae(wb)

    ordered_categories = []
    for cat in category_order_from_ae:
        if cat in cat_to_rows:
            ordered_categories.append(cat)
    for cat in categories_in_data:
        if cat not in ordered_categories:
            ordered_categories.append(cat)

    master_rows = []
    for cat in ordered_categories:
        dishes = cat_to_rows.get(cat, [])
        if not dishes:
            continue

        master_rows.append((True, cat, None, None, None))

        for _, name, weight, portions, gpp in dishes:
            master_rows.append(
                (
                    False,
                    str(name or ""),
                    weight if not skip_columns else None,
                    portions if not skip_columns else None,
                    gpp,
                )
            )

    return master_rows, total_food, total_liquid


def split_master_rows_to_slides(master_rows):
    max_rows_per_slide = int(
        (MAX_TABLE_HEIGHT_CM - ROW_HEIGHT_HEADER_CM) / ROW_HEIGHT_DATA_CM
    )
    slides = []
    i = 0
    n = len(master_rows)
    while i < n:
        slides.append(master_rows[i:i + max_rows_per_slide])
        i += max_rows_per_slide
    return slides


def set_table_style_no_grid(table):
    tbl = table._tbl
    tblPr_xml = """
    <a:tblPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
             xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
             xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
             firstRow="1" bandRow="1">
      <a:tableStyleId>{2D5ABB26-0587-4C30-8999-92F81FD0307C}</a:tableStyleId>
    </a:tblPr>
    """
    new_tblPr = parse_xml(tblPr_xml)
    existing_tblPr = tbl.xpath('./a:tblPr')
    if existing_tblPr:
        tbl.remove(existing_tblPr[0])
    tbl.insert(0, new_tblPr)


def create_slide_with_table(
    prs,
    header_text,
    bg_image_path: Path,
    slide_rows,
    skip_columns,
    is_last_slide,
    total_food_per_person,
    total_liquid_per_person,
):
    data_rows = len(slide_rows)
    extra_rows = 3 if is_last_slide else 0
    total_rows = 1 + data_rows + extra_rows

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if bg_image_path.exists():
        slide.shapes.add_picture(
            str(bg_image_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    tb = slide.shapes.add_textbox(Cm(1.5), Cm(1.0), prs.slide_width - Cm(3), Cm(1.8))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = header_text
    p.font.name = "Century Gothic"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    table_total_width = Cm(
        COL1_WIDTH_CM + COL2_WIDTH_CM + COL3_WIDTH_CM + COL4_WIDTH_CM
    )
    left = int((prs.slide_width - table_total_width) / 2)
    top = Cm(3.2)
    height = Cm(MAX_TABLE_HEIGHT_CM)

    shape = slide.shapes.add_table(total_rows, 4, left, top, table_total_width, height)
    table = shape.table

    set_table_style_no_grid(table)

    table.columns[0].width = Cm(COL1_WIDTH_CM)
    table.columns[1].width = Cm(COL2_WIDTH_CM)
    table.columns[2].width = Cm(COL3_WIDTH_CM)
    table.columns[3].width = Cm(COL4_WIDTH_CM)

    table.rows[0].height = Cm(ROW_HEIGHT_HEADER_CM)
    for r in range(1, total_rows):
        table.rows[r].height = Cm(ROW_HEIGHT_DATA_CM)

    for r in range(total_rows):
        for c in range(4):
            cell = table.cell(r, c)
            cell.fill.background()
            cell.margin_left = Cm(0)
            cell.margin_right = Cm(0)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT

    ws11 = getattr(prs, "_headers_ws", None)

    hdr_w = "Вес порции, грамм"
    hdr_p = "Кол-во порций"
    hdr_g = "Вес на одну персону, грамм"
    label_food = "Итого выход напитков на персону, мл"
    label_liquid = "Итого выход напитков на персону, мл"

    if ws11 is not None:
        hdr_w = str(ws11["A1"].value or hdr_w)
        hdr_p = str(ws11["B1"].value or hdr_p)
        hdr_g = str(ws11["C1"].value or hdr_g)
        label_food = str(ws11["A4"].value or label_food)
        label_liquid = str(ws11["A5"].value or label_liquid)

    cell = table.cell(0, 0)
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Наименования блюд"
    p.font.bold = True
    p.font.name = "Century Gothic"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    if skip_columns:
        table.cell(0, 1).text = ""
        table.cell(0, 2).text = ""
    else:
        cell = table.cell(0, 1)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = hdr_w
        p.font.bold = True
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        cell = table.cell(0, 2)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = hdr_p
        p.font.bold = True
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    cell = table.cell(0, 3)
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = hdr_g
    p.font.bold = True
    p.font.name = "Century Gothic"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    for idx, row in enumerate(slide_rows, start=0):
        is_category, text, weight, portions, gpp = row
        row_idx = 1 + idx

        cell = table.cell(row_idx, 0)
        if is_category:
            cell.margin_left = Cm(0)
        else:
            cell.margin_left = Cm(0.6)

        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text or ""
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.bold = bool(is_category)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

        if is_category:
            w_text = ""
            q_text = ""
            g_text = ""
        else:
            w_text = "" if (skip_columns or weight is None) else str(weight)
            q_text = "" if (skip_columns or portions is None) else str(portions)
            if isinstance(gpp, (int, float)):
                g_text = f"{float(gpp):.2f}".replace(".", ",")
            else:
                g_text = "" if gpp is None else str(gpp)

        cell = table.cell(row_idx, 1)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = w_text
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        cell = table.cell(row_idx, 2)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = q_text
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        cell = table.cell(row_idx, 3)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = g_text
        p.font.name = "Century Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    if is_last_slide:
        total_rows = len(table.rows)
        if total_rows >= 4:
            row_blank = total_rows - 3
            row_food = total_rows - 2
            row_liquid = total_rows - 1

            for c in range(4):
                table.cell(row_blank, c).text = ""

            cell = table.cell(row_food, 0)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label_food + ":"
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            tf.margin_left = Cm(0)

            cell = table.cell(row_food, 3)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{total_food_per_person:.2f}".replace(".", ",")
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

            cell = table.cell(row_liquid, 0)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label_liquid + ":"
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            tf.margin_left = Cm(0)

            cell = table.cell(row_liquid, 3)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{total_liquid_per_person:.2f}".replace(".", ",")
            p.font.name = "Century Gothic"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER


def process_sheet(
    wb,
    ws,
    sheet_index: int,
    prs: Presentation,
    bg_image_path: Path,
    skip_columns: bool,
):
    rows_raw = collect_rows_for_sheet(wb, ws, skip_columns)
    if not rows_raw:
        return

    master_rows, total_food_per_person, total_liquid_per_person = build_master_rows_and_totals(
        wb, rows_raw, skip_columns
    )
    if not master_rows:
        return

    slides = split_master_rows_to_slides(master_rows)
    if not slides:
        return

    last_rows_count = len(slides[-1])
    main_height_last = ROW_HEIGHT_HEADER_CM + last_rows_count * ROW_HEIGHT_DATA_CM
    can_place_totals_on_last = main_height_last <= MAIN_HEIGHT_LAST_CM

    header_text = get_header_text(wb, ws, sheet_index)

    for idx, slide_rows in enumerate(slides, start=1):
        is_last = (idx == len(slides)) and can_place_totals_on_last
        create_slide_with_table(
            prs,
            header_text,
            bg_image_path,
            slide_rows,
            skip_columns,
            is_last,
            total_food_per_person,
            total_liquid_per_person,
        )

    if not can_place_totals_on_last:
        create_slide_with_table(
            prs,
            header_text,
            bg_image_path,
            [],
            skip_columns,
            True,
            total_food_per_person,
            total_liquid_per_person,
        )


def build_presentation(excel_path: Path, bg_image_path: Path, out_path: Path) -> Path:
    wb = load_book(excel_path)
    prs = Presentation()
    prs._headers_ws = get_headers_sheet(wb)  # type: ignore
    skip_columns = get_skip_columns_flag(wb)

    for idx, ws in enumerate(wb.worksheets, start=1):
        if 3 <= idx <= 8:
            process_sheet(wb, ws, idx, prs, bg_image_path, skip_columns)

    prs.save(out_path)
    return out_path
